"""This python script generates the PPTx deck that can be leveraged for Sprint Reviews instead of requiring manual intervention.
It creates a slide for each epic and puts all the stories that were moved to In Review, Ready for Prod, or Done status. One additional note is 
that it checks the jira story history to ensure the item was not already reviewed in a prior discussion."""


import common_functions as cf
import pandas as pd
import numpy as np
from dotenv import load_dotenv
import os
import ast
from pptx import Presentation
from pptx.util import Inches, Pt
import math
from pptx.dml.color import RGBColor
from office365.sharepoint.client_context import ClientContext
from office365.runtime.auth.user_credential import UserCredential
import requests
from datetime import datetime
import json 

def add_to_gitignore(file_name, gitignore_path):
    """
    Dynamically add file to gitignore. This is needed since each sprint the file name varies.
    
    :param file_name: Description
    """
    ignore_pattern = file_name.strip()
    with open(gitignore_path, 'a') as f:
        f.write(f'\n{ignore_pattern}')
        print(f'Added {ignore_pattern} to {gitignore_path}')


def create_sprint_pptx(template, pptx_title, sprint_start, sprint_end, data):
    """
    Generates a powerpoint document to summarize all work that has been moved to in review, ready for prod, or done in the current sprint for discussions/presentations.
    
    :param template: pptx file that has the set template that you want to follow
    :param pptx_title: Text you want to be as the main text on your title slide (Recommended the sprint name)
    :param sprint_start: Start date of sprint used for subtext in title slide
    :param sprint_end: End date of sprint that is used for subtext in title slide
    :param data: Dataframe that houses all data needed to generate the report. Initaitve Name, Epic Name, Issue ID, Summary, Asignee all required in this to work properly
    
    :return: Saves a pptx file to repository
    """
    print('-'*60)
    print(f"Creating PPtx File")
    print('-'*60)    
    prs = Presentation(template)
    print('----Checking Template File Formats----')
    for i, layout in enumerate(prs.slide_layouts):
        print(f'Index: {i}, Name: {layout.name}')
    slide_layout = prs.slide_layouts[2]
    epic_df = data[['Epic ID', 'Epic Name', 'Epic Description', 'Jira Board']].drop_duplicates()
    slide = prs.slides.add_slide(prs.slide_layouts[0]) 
    title = slide.shapes.title 
    title.text = pptx_title 
    subtitle = slide.placeholders[1] 
    subtitle.text = str(sprint_start) + ' - ' + str(sprint_end)
    tb_left = Inches(0.5)  
    tb_top = Inches(1.5)
    tb_width = Inches(8)
    tb_height = Inches(1)
    print('Title Slide Created')
    initiative_names = data['Initiative'].unique().tolist()
    print(initiative_names)
    for initiative in range(len(initiative_names)):
        print(initiative_names[initiative])
        epic_names = data[data['Initiative'] == initiative_names[initiative]]['Epic Name'].unique()
        print(epic_names)
        for epic in range(len(epic_names)):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            text_frame = title.text_frame
            p1 = text_frame.paragraphs[0]
            p1.clear()
            run2 = p1.add_run()
            run2.font.size = Pt(20)
            run2.font.bold = True
            run2.font.color.rgb = RGBColor(65, 79, 100)  # Set font color to black
            run2.text = epic_df[epic_df['Epic Name'] == epic_names[epic]]['Jira Board'].values[0].split()[1] + ' | '
            run3 = p1.add_run()
            run3.font.size = Pt(20)
            run3.font.color.rgb = RGBColor(89, 171, 221)  # Set font color to black
            run3.font.bold = True
            run3.text = initiative_names[initiative]
            tb_left = Inches(0.5)  
            tb_top = Inches(1.5)
            tb_width = Inches(8)
            tb_height = Inches(1)
            textbox = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
            text_frame = textbox.text_frame
            # # Add a paragraph to the text frame
            placeholder = slide.placeholders[1]  # Use the second placeholder for the text
            p1 = placeholder.text_frame.paragraphs[0]
            run1 = p1.add_run()
            run1.text = data[data['Epic Name'] == epic_names[epic]]['Epic Name'].values[0]
            run1.font.bold = True
            p = placeholder.text_frame.add_paragraph()
            text_frame.word_wrap = True
            # Handle NaN values and ensure string conversion for PowerPoint
            epic_description = epic_df[epic_df['Epic Name'] == epic_names[epic]]['Epic Description'].values[0]
            p.text = str(epic_description) if pd.notna(epic_description) and epic_description is not None else 'No Epic Description provided'
            # print(epic_names[x])
            # print( items[(items['Initiative Name'] == initiative_names[i]) & (items['Epic ID'] == epic_names[x])][['Summary', 'Assignee', 'Status', 'Environment']])
            # Add a table to the slide
            # print(epic_names[x])
                # print(items[(items['Initiative Name'] == initiative_names[i]) & (items['Epic'] == epic_names[x])].iloc[z][['Summary', 'Assignee', 'Status', 'Environment']])
            rows, cols = data[(data['Initiative'] == initiative_names[initiative]) & (data['Epic Name'] == epic_names[epic])][['Summary', 'Assignee', 'Status', 'Environment', 'Fix Version']].shape
            table_top = tb_top + tb_height + Inches(1.2)  # 0.2 inch gap
            table_left = tb_left + Inches(0.3)  # 0.1 inch gap
            table_width = tb_width
            table_height = Inches(0.8 + 0.3*rows) 
            table = slide.shapes.add_table(rows + 1, cols, table_left, table_top, table_width, table_height).table
            for col_idx, col_name in enumerate(['Fix Version','Summary', 'Assignee', 'Status', 'Environment']):
                table.cell(0, col_idx).text = col_name
                # Set font size for column names
                for paragraph in table.cell(0, col_idx).text_frame.paragraphs:
                    for run in paragraph.runs:
                        cell = table.cell(0, col_idx)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(65, 79, 100)
                        run.font.size = Pt(12)
            # Fill the table with data
            for row_idx in range(rows): 
                for col_idx in range(cols):
                    cell_value = str(data[(data['Initiative'] == initiative_names[initiative]) & (data['Epic Name'] == epic_names[epic])][['Fix Version','Summary', 'Assignee', 'Status', 'Environment']].iat[row_idx, col_idx])
                    table.cell(row_idx + 1, col_idx).text = cell_value
                    # Set font size for cell values
                    for paragraph in table.cell(row_idx + 1, col_idx).text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            cell = table.cell(row_idx + 1, col_idx)
                            cell.fill.solid()       
                            cell.fill.fore_color.rgb = RGBColor(89, 171, 221)  # Set cell background color to white
                            print('adding items')

    pptx_file = f'{pptx_title} Showcase.pptx'.replace(' ', '_')
    prs.save(pptx_file)
    return pptx_file

def find_doclibid(json_doclibdata, doclib_name):
    return [p for p in json_doclibdata if p['name']==doclib_name][0]['id']

def publish_to_sharepoint(file, authority, client_id, client_secret, scope, tentant_domain, site_name, doclib):
    """
    Allows you to publish a file to selected url by passing in the following information. In order for this function to work, the user needs to be configured to have access to publish through Microsoft Graph.
    
    :param file: the file location path that you are looking to publish to the site
    :param url: the url of the sharepoint domain you are looking to publish to
    :param user: username that is used as credentials to publish the file
    :param password: password of user that is used as credentials to publish the file
    :param location: The file path within your sharepoint site that you would like to publish the file to
    """
    app_settings = {
        'authority': f'{authority}',
        'client_id': f'{client_id}',
        'client_secret': f'{client_secret}',
        'scope': [f'{scope}'],
        'tenantdomain': f'{tentant_domain}',
        'sitename': f'{site_name}',
        'doclib': f'{doclib}',
    }       

    token_url = app_settings['authority'] + '/oauth2/v2.0/token'
    token_json_body = dict(
        grant_type="client_credentials",
        scope=app_settings["scope"],
        client_id=app_settings["client_id"],
        client_secret=app_settings["client_secret"]
    )
    token_result = requests.post(token_url,token_json_body)
    token_result_data = token_result.json()
    graph_token_header = {'Authorization': 'Bearer ' + token_result_data['access_token']}
    get_site_id_url = 'https://graph.microsoft.com/v1.0/sites/'+ app_settings["tenantdomain"] + ':/sites/' + app_settings["sitename"] + '?$select=id,displayName'

    # Calling graph using the access token
    graph_site_data = requests.get(  # Use token to call downstream service
        get_site_id_url,
        headers = graph_token_header 
    ,).json()

    graph_get_doc_lib = 'https://graph.microsoft.com/v1.0/sites/' + graph_site_data['id'] + '/drives'
    graph_doc_lib_data = requests.get(  # Use token to call downstream service
    graph_get_doc_lib,
    headers = graph_token_header 
    ,).json()
    graph_doc_lib_id = find_doclibid(graph_doc_lib_data['value'], app_settings["doclib"])

    file_path = file
    file_name = os.path.basename(file_path)
    file_size = os.path.getsize(file_path)

    #Upload very large file to folder Vault
    url = 'https://graph.microsoft.com/v1.0/drives/' + graph_doc_lib_id + '/root:/General/' + file_name + ':/createUploadSession?@microsoft.graph.conflictBehavior=replace'
    url = json.loads(requests.post(url, headers=graph_token_header).text)
    url = url['uploadUrl']
    
    chunk_size = 320*1024*10 # Has to be multiple of 320 kb
    no_of_uploads = file_size//chunk_size
    content_range_start = 0
    if file_size < chunk_size :
        content_range_end = file_size
    else :
        content_range_end = chunk_size - 1
    
    data = open(file_path, 'rb')
    while data.tell() < file_size:
        if ((file_size - data.tell()) <= chunk_size):
            content_range_end = file_size -1
            headers = {'Content-Range' : 'bytes '+ str(content_range_start)+ '-' +str(content_range_end)+'/'+str(file_size)}
            content = data.read(chunk_size)
            response = json.loads(requests.put(url, headers=headers, data = content).text)
        else:
            headers = {'Content-Range' : 'bytes '+ str(content_range_start)+ '-' +str(content_range_end)+'/'+str(file_size)}
            content = data.read(chunk_size)
            response = json.loads(requests.put(url, headers=headers, data = content).text)
            content_range_start = data.tell()
            content_range_end = data.tell() + chunk_size - 1
    data.close()
    response2 = requests.delete(url)

    if os.path.exists(file_path):
        os.remove(file_path)


if __name__ == '__main__':
    tableau_user = os.environ['TABLEAU_USER']
    tableau_password = os.environ['TABLEAU_PASSWORD']
    tableau_site = os.environ['TABLEAU_SITE']
    tableau_server = os.environ['TABLEAU_SERVER']
    jira_url = os.environ['JIRA_URL']
    jira_user = os.environ['JIRA_USER']
    api_token = os.environ['JIRA_TOKEN']
    custom_fields = ast.literal_eval(os.environ['CUSTOM_FIELDS'])
    sharepoint_domain = os.environ['SHAREPOINT_DOMAIN']
    sharepoint_secret = os.environ['SHAREPOINT_CLIENT_SECRET']
    sharepoint_client_id = os.environ['SHAREPOINT_CLIENT_ID']
    sharepoint_site_name = os.environ['SHAREPOINT_SITE_NAME']
    sharepoint_doc_lib = os.environ['SHAREPOINT_DOC_LIB']
    sharepoint_authority = os.environ['SHAREPOINT_AUTHORITY']
    sharepoint_scope = os.environ['SHAREPOINT_SCOPE']

    # Get all the issues and information in the Sprint
    # Project details were not moved to .env file so masked with some random numbers and project info

    isg_project = cf.JiraProject('Project 1', '1111', '11111', jira_url, jira_user, api_token, custom_fields)
    asg_project = cf.JiraProject('Project 2', '2222', '22222', jira_url, jira_user, api_token, custom_fields)
    project_list = [isg_project, asg_project]
    df_columns = ['Issue ID', 'Issue Type', 'Summary', 'Description', 'Status',
       'Assignee', 'Reporter', 'Created Date', 'Updated Date', 'Epic Name',
       'Epic ID', 'Story Points', 'Jira Board', 'Fix Version', 'Sprint',
       'Environment', 'Epic Description', 'Issue Summary', 'Initiative ID',
       'Initiative', 'Initiative Type', 'Initiative Status']
    sprint_df = pd.DataFrame(columns=df_columns)
    for project in project_list:
        # Loop through each project we need for our Sprint Deck
        print(f'JIRA Project: {project.project}')
        # First Retrieve the Current Sprint detail that we want to create the deck for
        sprints = project.get_sprint_details()
        sprint_name = sprints[sprints['state'] == 'active']['name'].iloc[0]
        sprint_start = pd.to_datetime(sprints[sprints['state'] == 'active']['startDate'].iloc[0]).date()
        sprint_end = pd.to_datetime(sprints[sprints['state'] == 'active']['endDate'].iloc[0]).date()

        # Fetch Issues, Epic Details, and Initiatives for use in Pptx File
        issue_df = project.get_issues_in_sprint(sprint_id=sprint_name, type=['Bug', 'Story'])
        initiative_df = project.get_parent_details(issue_df['Parent ID'], maxResults=False)
        epic_list = issue_df['Parent ID'].unique().tolist()
        if None in epic_list:
            epic_list.remove(None)
        epic_df = project.get_all_issues(issue_ids=epic_list, type=['Epic'], cust_fields=None)

        # Merge Issues with Epic as Epic Description is not returned without the get all issues function
        issue_join_df = pd.merge(issue_df, epic_df[['Issue ID', 'Description']], how='left', left_on='Parent ID', right_on='Issue ID', suffixes=('_Child', '_Parent'))
        issue_join_df = issue_join_df.drop(columns=(['Issue ID_Parent']))
        issue_join_df = issue_join_df.rename(columns={'Issue ID_Child': 'Issue ID', 'Description_Parent': 'Epic Description',
                                                    'Description_Child': 'Description'})
        # Merge initiative details into issues & remove any noise of column naming
        prepped_df = pd.merge(issue_join_df, initiative_df, how='left', left_on='Parent ID', right_on='Issue ID', suffixes=('_Child', '_Parent'))
        prepped_df = prepped_df.drop(columns=(['Issue ID_Parent', 'Issue Assignee', 'Issue Type_Parent']))
        prepped_df = prepped_df.rename(columns={'Issue ID_Child': 'Issue ID', 'Issue Type_Child': 'Issue Type',
                                                            'Parent ID_Child': 'Epic ID', 'Parent': 'Epic Name',
                                                            'Parent ID_Parent': "Initiative ID", 'Parent Summary': 'Initiative',
                                                            'Parent Type': 'Initiative Type', 'Parent Status': 'Initiative Status'})
        # Get audit history and Create an exclusion list of stories that were moved to In Review, Ready for Prod, or Done in a prior sprint (Some statuses are not defined as done in JIRA)
        # Also remove any mapping or QA work as we do not demo those
        audit_issues = project.get_audit_log(issue_df['Issue ID'].unique().tolist())
        exclusion_list = audit_issues[(audit_issues['Field'] == 'status') & (audit_issues['ChangeDate'] < sprint_start) & (audit_issues['ToStatus'].isin(['Done', 'In Review', 'Ready for Prod']))]['Issue ID'].unique().tolist()    
        print(f'Total # of stories excluded: {len(exclusion_list)}')

        #We had some individuals that were doing testing or data mapping and we didn't feel the need to showcase
        final_output = prepped_df[(~prepped_df['Issue ID'].isin(exclusion_list)) & (~prepped_df['Assignee'].isin(['User 1','User 2']))
                                                                                    & (prepped_df['Status'].isin(['Done', 'In Review', 'Ready for Prod']))]
        sprint_df = pd.concat([sprint_df, final_output])

    # Use robust path to find PowerPoint template (works locally and in GitHub Actions)
    template_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'Sprint_showcase.pptx')
    pptx_file = create_sprint_pptx(template=template_path, pptx_title=sprint_name, data=sprint_df, sprint_start=sprint_start, sprint_end=sprint_end)
    publish_to_sharepoint(file=pptx_file, authority=sharepoint_authority, client_id=sharepoint_client_id, client_secret=sharepoint_secret, scope=sharepoint_scope, tentant_domain=sharepoint_domain, site_name=sharepoint_site_name, doclib=sharepoint_doc_lib)
    
